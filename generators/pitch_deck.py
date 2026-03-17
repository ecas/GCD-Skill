"""
Google Cloud Pitch Deck Generator
==================================
Generates branded Google Cloud sales pitch decks as .pptx files.

Usage (standalone):
    python pitch_deck.py --config deck_config.json --output output.pptx

Usage (as library):
    from generators.pitch_deck import generate_pitch_deck
    generate_pitch_deck(config)

Requires: python-pptx>=0.6.21
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor


# ---------------------------------------------------------------------------
# Brand constants
# ---------------------------------------------------------------------------

BLUE = RGBColor(0x42, 0x85, 0xF4)          # Google Blue
GREEN = RGBColor(0x34, 0xA8, 0x53)         # Google Green
YELLOW = RGBColor(0xFB, 0xBC, 0x04)        # Google Yellow
RED = RGBColor(0xEA, 0x43, 0x35)           # Google Red
DARK_BLUE = RGBColor(0x1A, 0x23, 0x7E)     # Dark navy (title bg)
LIGHT_BLUE_BG = RGBColor(0xE8, 0xF0, 0xFE) # Table alt row / accents
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x20, 0x20, 0x20)
MEDIUM_GREY = RGBColor(0x75, 0x75, 0x75)
LIGHT_GREY = RGBColor(0xF8, 0xF9, 0xFA)
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)        # Timeline phase 3
HEADER_BAR = BLUE

# Slide dimensions — 16:9 widescreen (standard)
SLIDE_WIDTH = Cm(33.867)
SLIDE_HEIGHT = Cm(19.05)

FONT_FAMILY = "Calibri"  # Google Sans is ideal but not embedded; Calibri is safe cross-platform

# Layout zones (in Cm)
HEADER_TOP = Cm(0)
HEADER_HEIGHT = Cm(2.0)
CONTENT_TOP = Cm(2.4)
CONTENT_LEFT = Cm(1.2)
CONTENT_WIDTH = SLIDE_WIDTH - Cm(2.4)
CONTENT_HEIGHT = SLIDE_HEIGHT - Cm(3.2)
FOOTER_TOP = SLIDE_HEIGHT - Cm(0.7)
FOOTER_HEIGHT = Cm(0.7)


# ---------------------------------------------------------------------------
# Low-level drawing helpers
# ---------------------------------------------------------------------------

def _rgb(r: int, g: int, b: int) -> RGBColor:
    return RGBColor(r, g, b)


def _add_rect(
    slide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    fill_color: RGBColor,
    line_color: RGBColor | None = None,
    line_width_pt: float = 0,
):
    """Add a filled rectangle shape to the slide."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def _add_textbox(
    slide,
    text: str,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    font_size: int = 14,
    bold: bool = False,
    color: RGBColor = DARK_TEXT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_family: str = FONT_FAMILY,
    wrap: bool = True,
    italic: bool = False,
):
    """Add a plain textbox to the slide."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_family
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def _add_text_to_shape(
    shape,
    text: str,
    font_size: int = 14,
    bold: bool = False,
    color: RGBColor = WHITE,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_family: str = FONT_FAMILY,
):
    """Set text on an existing shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_family
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_footer(slide, slide_number: int | None, confidential: bool = True):
    """Add footer with slide number and confidential label."""
    footer_text = "Confidential — Google Confidential" if confidential else ""
    _add_textbox(
        slide,
        footer_text,
        left=Cm(1.0),
        top=FOOTER_TOP,
        width=Cm(20),
        height=FOOTER_HEIGHT,
        font_size=8,
        color=MEDIUM_GREY,
        italic=True,
    )
    if slide_number is not None:
        _add_textbox(
            slide,
            str(slide_number),
            left=SLIDE_WIDTH - Cm(2),
            top=FOOTER_TOP,
            width=Cm(1.5),
            height=FOOTER_HEIGHT,
            font_size=8,
            color=MEDIUM_GREY,
            align=PP_ALIGN.RIGHT,
        )


def _add_header_bar(slide, title: str):
    """Add the blue header bar and title text used on content slides."""
    bar = _add_rect(
        slide,
        left=Cm(0),
        top=HEADER_TOP,
        width=SLIDE_WIDTH,
        height=HEADER_HEIGHT,
        fill_color=HEADER_BAR,
    )
    _add_text_to_shape(
        bar,
        title,
        font_size=20,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.LEFT,
    )
    bar.text_frame.paragraphs[0].runs[0].font.name = FONT_FAMILY
    # Indent text within bar
    from pptx.util import Pt as _Pt
    bar.text_frame.margin_left = Cm(0.5)
    bar.text_frame.margin_top = Cm(0.2)


def _add_table(
    slide,
    headers: list[str],
    rows: list[list[str]],
    left: Emu,
    top: Emu,
    width: Emu,
    col_widths: list[float] | None = None,
    header_color: RGBColor = HEADER_BAR,
    alt_row_color: RGBColor = LIGHT_BLUE_BG,
    font_size: int = 12,
):
    """Add a formatted table to the slide."""
    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header
    row_height = Cm(0.85)
    table_height = row_height * n_rows

    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_height).table

    if col_widths:
        total = sum(col_widths)
        for idx, cw in enumerate(col_widths):
            table.columns[idx].width = int(width * cw / total)

    # Header row
    for col_idx, header_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = header_text
        run.font.name = FONT_FAMILY
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = WHITE

    # Data rows
    for row_idx, row_data in enumerate(rows):
        bg = alt_row_color if row_idx % 2 == 1 else WHITE
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(cell_text)
            run.font.name = FONT_FAMILY
            run.font.size = Pt(font_size)
            run.font.color.rgb = DARK_TEXT

    return table


# ---------------------------------------------------------------------------
# Slide renderers
# ---------------------------------------------------------------------------

def _render_title(slide, data: dict[str, Any], config: dict[str, Any]):
    """
    Title slide: dark navy background, large headline, subtitle, company/date/presenter.
    """
    # Full-bleed dark background
    _add_rect(slide, Cm(0), Cm(0), SLIDE_WIDTH, SLIDE_HEIGHT, DARK_BLUE)

    # Google Cloud color bar at bottom (4 stripes)
    stripe_w = SLIDE_WIDTH / 4
    for idx, color in enumerate([BLUE, RED, YELLOW, GREEN]):
        _add_rect(
            slide,
            left=int(stripe_w * idx),
            top=SLIDE_HEIGHT - Cm(0.5),
            width=int(stripe_w),
            height=Cm(0.5),
            fill_color=color,
        )

    # "Google Cloud" label (logo placeholder)
    _add_textbox(
        slide,
        "Google Cloud",
        left=Cm(1.5),
        top=Cm(1.2),
        width=Cm(10),
        height=Cm(1.0),
        font_size=18,
        bold=True,
        color=WHITE,
    )

    headline = data.get("headline", "")
    subtitle = data.get("subtitle", "")
    company = config.get("company_name", "")
    presenter = config.get("presenter", "")
    date = config.get("date", "")

    # Main headline
    _add_textbox(
        slide,
        headline,
        left=Cm(1.5),
        top=Cm(5.5),
        width=SLIDE_WIDTH - Cm(3),
        height=Cm(3.5),
        font_size=40,
        bold=True,
        color=WHITE,
    )

    # Subtitle
    if subtitle:
        _add_textbox(
            slide,
            subtitle,
            left=Cm(1.5),
            top=Cm(9.5),
            width=SLIDE_WIDTH - Cm(3),
            height=Cm(1.8),
            font_size=22,
            bold=False,
            color=_rgb(0xBB, 0xDE, 0xFB),
        )

    # Company / presenter / date strip
    meta = f"{company}   |   {presenter}   |   {date}"
    _add_textbox(
        slide,
        meta,
        left=Cm(1.5),
        top=SLIDE_HEIGHT - Cm(2.0),
        width=SLIDE_WIDTH - Cm(3),
        height=Cm(0.8),
        font_size=11,
        color=_rgb(0x90, 0xCA, 0xF9),
    )

    slide.shapes.title  # suppress default title placeholder warnings

    slide.notes_slide.notes_text_frame.text = (
        f"Title slide for {company}. Presenter: {presenter}. "
        "Introduce yourself and the agenda briefly. Confirm the meeting objective with the audience."
    )


def _render_agenda(slide, data: dict[str, Any], slide_number: int):
    """Agenda slide with numbered items and optional timing column."""
    _add_header_bar(slide, "Agenda")
    _add_footer(slide, slide_number)

    items = data.get("items", [])
    timings = data.get("timing", [])

    left = CONTENT_LEFT
    top = CONTENT_TOP + Cm(0.3)
    item_height = Cm(1.5)

    for idx, item in enumerate(items):
        # Number circle
        circle = _add_rect(
            slide,
            left=left,
            top=top + item_height * idx,
            width=Cm(0.9),
            height=Cm(0.9),
            fill_color=BLUE,
        )
        _add_text_to_shape(circle, str(idx + 1), font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Item text
        _add_textbox(
            slide,
            item,
            left=left + Cm(1.1),
            top=top + item_height * idx,
            width=Cm(22),
            height=Cm(0.9),
            font_size=16,
            color=DARK_TEXT,
        )

        # Timing badge
        if idx < len(timings):
            badge = _add_rect(
                slide,
                left=SLIDE_WIDTH - Cm(3.5),
                top=top + item_height * idx,
                width=Cm(2.8),
                height=Cm(0.75),
                fill_color=LIGHT_BLUE_BG,
            )
            _add_text_to_shape(
                badge,
                timings[idx],
                font_size=11,
                bold=False,
                color=BLUE,
                align=PP_ALIGN.CENTER,
            )

    slide.notes_slide.notes_text_frame.text = (
        "Walk through the agenda. Confirm any topics the audience wants to add or skip. "
        "Set expectations: interactive session, questions welcome throughout."
    )


def _render_pain_points(slide, data: dict[str, Any], slide_number: int):
    """Pain points slide with red-accented numbered items and source attribution."""
    headline = data.get("headline", "Challenges")
    _add_header_bar(slide, headline)
    _add_footer(slide, slide_number)

    points = data.get("points", [])

    top = CONTENT_TOP + Cm(0.5)
    spacing = (CONTENT_HEIGHT - Cm(0.5)) / max(len(points), 1)

    for idx, point in enumerate(points):
        pain = point.get("pain", "")
        source = point.get("source", "")
        item_top = top + spacing * idx

        # Red accent circle with number
        circle = _add_rect(
            slide,
            left=CONTENT_LEFT,
            top=item_top,
            width=Cm(1.0),
            height=Cm(1.0),
            fill_color=RED,
        )
        _add_text_to_shape(
            circle,
            str(idx + 1),
            font_size=13,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )

        # Pain text
        _add_textbox(
            slide,
            pain,
            left=CONTENT_LEFT + Cm(1.3),
            top=item_top,
            width=CONTENT_WIDTH - Cm(5),
            height=Cm(0.9),
            font_size=16,
            bold=True,
            color=DARK_TEXT,
        )

        # Source attribution
        if source:
            _add_textbox(
                slide,
                f"Source: {source}",
                left=CONTENT_LEFT + Cm(1.3),
                top=item_top + Cm(0.85),
                width=CONTENT_WIDTH - Cm(5),
                height=Cm(0.55),
                font_size=10,
                color=MEDIUM_GREY,
                italic=True,
            )

    slide.notes_slide.notes_text_frame.text = (
        "These pain points should resonate with the audience. "
        "Pause after each one and ask: 'Does this match what your team is experiencing?' "
        "Use their language — don't read from the slide."
    )


def _render_solution_map(slide, data: dict[str, Any], slide_number: int):
    """Solution map: 3-column table mapping challenges to solutions and outcomes."""
    headline = data.get("headline", "Our Solution")
    _add_header_bar(slide, headline)
    _add_footer(slide, slide_number)

    rows_data = data.get("rows", [])
    table_rows = [[r.get("challenge", ""), r.get("solution", ""), r.get("outcome", "")] for r in rows_data]

    _add_table(
        slide,
        headers=["Challenge", "Google Cloud Solution", "Business Outcome"],
        rows=table_rows,
        left=CONTENT_LEFT,
        top=CONTENT_TOP + Cm(0.5),
        width=CONTENT_WIDTH,
        col_widths=[3, 3, 4],
        font_size=13,
    )

    slide.notes_slide.notes_text_frame.text = (
        "This is the core value proposition. Walk each row: confirm the challenge is real for them, "
        "name the specific Google Cloud product, and quantify the outcome if possible."
    )


def _render_comparison_table(slide, data: dict[str, Any], slide_number: int):
    """Before/after comparison table."""
    headline = data.get("headline", "Current State vs Google Cloud")
    _add_header_bar(slide, headline)
    _add_footer(slide, slide_number)

    headers = data.get("headers", ["Dimension", "Today", "With Google Cloud"])
    rows = data.get("rows", [])

    _add_table(
        slide,
        headers=headers,
        rows=rows,
        left=CONTENT_LEFT,
        top=CONTENT_TOP + Cm(0.5),
        width=CONTENT_WIDTH,
        font_size=13,
    )

    slide.notes_slide.notes_text_frame.text = (
        "Emphasise the contrast between today's reality and what Google Cloud enables. "
        "Ask: 'Which of these dimensions is most painful right now?'"
    )


def _render_customer_proof(slide, data: dict[str, Any], slide_number: int):
    """Customer proof points — industry reference wins."""
    headline = data.get("headline", "Customer Proof Points")
    _add_header_bar(slide, headline)
    _add_footer(slide, slide_number)

    stories = data.get("stories", [])

    left = CONTENT_LEFT
    top = CONTENT_TOP + Cm(0.5)
    card_width = (CONTENT_WIDTH - Cm(0.5) * (len(stories) - 1)) / max(len(stories), 1)

    for idx, story in enumerate(stories):
        card_left = left + (card_width + Cm(0.5)) * idx
        company = story.get("company", "")
        result = story.get("result", "")
        detail = story.get("detail", "")

        # Card background
        card = _add_rect(
            slide,
            left=card_left,
            top=top,
            width=card_width,
            height=Cm(8),
            fill_color=LIGHT_BLUE_BG,
        )

        # Company name
        _add_textbox(
            slide,
            company,
            left=card_left + Cm(0.3),
            top=top + Cm(0.4),
            width=card_width - Cm(0.6),
            height=Cm(1.0),
            font_size=20,
            bold=True,
            color=DARK_BLUE,
        )

        # Divider bar
        _add_rect(
            slide,
            left=card_left + Cm(0.3),
            top=top + Cm(1.5),
            width=card_width - Cm(0.6),
            height=Cm(0.06),
            fill_color=BLUE,
        )

        # Result headline
        _add_textbox(
            slide,
            result,
            left=card_left + Cm(0.3),
            top=top + Cm(1.8),
            width=card_width - Cm(0.6),
            height=Cm(2.5),
            font_size=16,
            bold=True,
            color=GREEN,
        )

        # Detail / method
        _add_textbox(
            slide,
            detail,
            left=card_left + Cm(0.3),
            top=top + Cm(4.5),
            width=card_width - Cm(0.6),
            height=Cm(2.5),
            font_size=12,
            color=MEDIUM_GREY,
        )

    slide.notes_slide.notes_text_frame.text = (
        "Use these references to build credibility. If you can share the full case study, do so. "
        "Offer to connect the customer with these reference accounts if they want a peer conversation."
    )


def _render_roi_snapshot(slide, data: dict[str, Any], slide_number: int):
    """ROI snapshot: large metric numbers with supporting labels."""
    headline = data.get("headline", "The Financial Case")
    _add_header_bar(slide, headline)
    _add_footer(slide, slide_number)

    metrics = data.get("metrics", [])

    left = CONTENT_LEFT
    top = CONTENT_TOP + Cm(0.8)
    card_width = (CONTENT_WIDTH - Cm(0.5) * (len(metrics) - 1)) / max(len(metrics), 1)
    accent_colors = [GREEN, BLUE, YELLOW]

    for idx, metric in enumerate(metrics):
        card_left = left + (card_width + Cm(0.5)) * idx
        label = metric.get("label", "")
        value = metric.get("value", "")
        detail = metric.get("detail", "")
        color = accent_colors[idx % len(accent_colors)]

        # Accent top bar
        _add_rect(
            slide,
            left=card_left,
            top=top,
            width=card_width,
            height=Cm(0.2),
            fill_color=color,
        )

        # Value — large
        _add_textbox(
            slide,
            value,
            left=card_left,
            top=top + Cm(0.4),
            width=card_width,
            height=Cm(3.0),
            font_size=48,
            bold=True,
            color=color,
            align=PP_ALIGN.CENTER,
        )

        # Label
        _add_textbox(
            slide,
            label,
            left=card_left,
            top=top + Cm(3.5),
            width=card_width,
            height=Cm(1.2),
            font_size=14,
            bold=True,
            color=DARK_TEXT,
            align=PP_ALIGN.CENTER,
        )

        # Detail
        _add_textbox(
            slide,
            detail,
            left=card_left,
            top=top + Cm(4.8),
            width=card_width,
            height=Cm(1.0),
            font_size=11,
            color=MEDIUM_GREY,
            align=PP_ALIGN.CENTER,
        )

    slide.notes_slide.notes_text_frame.text = (
        "Walk through each metric. Where possible, right-size these numbers to the customer's scale. "
        "Ask: 'Which of these would unlock your business case?' Do not present all three as equal — "
        "prioritise based on what you learned in discovery."
    )


def _render_timeline(slide, data: dict[str, Any], slide_number: int):
    """Horizontal phased timeline."""
    headline = data.get("headline", "Implementation Timeline")
    _add_header_bar(slide, headline)
    _add_footer(slide, slide_number)

    phases = data.get("phases", [])
    phase_colors = [GREEN, BLUE, PURPLE]

    left = CONTENT_LEFT
    top = CONTENT_TOP + Cm(0.5)
    phase_width = (CONTENT_WIDTH - Cm(0.3) * (len(phases) - 1)) / max(len(phases), 1)

    for idx, phase in enumerate(phases):
        phase_left = left + (phase_width + Cm(0.3)) * idx
        color = phase_colors[idx % len(phase_colors)]
        name = phase.get("name", f"Phase {idx + 1}")
        duration = phase.get("duration", "")
        items = phase.get("items", [])

        # Phase header block
        header_block = _add_rect(
            slide,
            left=phase_left,
            top=top,
            width=phase_width,
            height=Cm(1.4),
            fill_color=color,
        )

        # Phase name + duration in header
        _add_text_to_shape(
            header_block,
            f"{name}",
            font_size=14,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )

        # Duration below header
        _add_textbox(
            slide,
            duration,
            left=phase_left,
            top=top + Cm(1.5),
            width=phase_width,
            height=Cm(0.6),
            font_size=11,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

        # Items
        for item_idx, item_text in enumerate(items):
            _add_textbox(
                slide,
                f"  {item_text}",
                left=phase_left + Cm(0.2),
                top=top + Cm(2.3) + Cm(0.75) * item_idx,
                width=phase_width - Cm(0.4),
                height=Cm(0.7),
                font_size=12,
                color=DARK_TEXT,
            )

        # Arrow connector (not between last phase)
        if idx < len(phases) - 1:
            arrow_left = phase_left + phase_width
            _add_rect(
                slide,
                left=arrow_left,
                top=top + Cm(0.5),
                width=Cm(0.3),
                height=Cm(0.4),
                fill_color=MEDIUM_GREY,
            )

    slide.notes_slide.notes_text_frame.text = (
        "Walk the timeline left to right. Emphasise that Phase 1 is low-risk and high-value — "
        "a contained pilot that proves ROI before committing to wider rollout. "
        "Ask: 'Does this timeline match your internal planning horizon?'"
    )


def _render_next_steps(slide, data: dict[str, Any], slide_number: int):
    """Next steps action table."""
    headline = data.get("headline", "Next Steps")
    _add_header_bar(slide, headline)
    _add_footer(slide, slide_number)

    actions = data.get("actions", [])
    table_rows = [[a.get("action", ""), a.get("owner", ""), a.get("date", "")] for a in actions]

    _add_table(
        slide,
        headers=["Action", "Owner", "Target Date"],
        rows=table_rows,
        left=CONTENT_LEFT,
        top=CONTENT_TOP + Cm(0.5),
        width=CONTENT_WIDTH,
        col_widths=[6, 2, 2],
        font_size=13,
    )

    slide.notes_slide.notes_text_frame.text = (
        "Close by confirming each action item, owner, and date out loud. "
        "Ask for agreement before leaving the room. Offer to send the summary by end of day. "
        "The goal is a committed next meeting with a technical audience."
    )


def _render_text_slide(slide, data: dict[str, Any], slide_number: int):
    """Generic text slide with headline and body."""
    headline = data.get("headline", "")
    body = data.get("body", "")

    _add_header_bar(slide, headline)
    _add_footer(slide, slide_number)

    _add_textbox(
        slide,
        body,
        left=CONTENT_LEFT,
        top=CONTENT_TOP + Cm(0.5),
        width=CONTENT_WIDTH,
        height=CONTENT_HEIGHT - Cm(0.5),
        font_size=16,
        color=DARK_TEXT,
    )

    slide.notes_slide.notes_text_frame.text = (
        f"Slide: {headline}. Expand on the key points verbally. Keep this slide as a prompt — "
        "do not read it word for word."
    )


def _render_architecture(slide, data: dict[str, Any], slide_number: int):
    """
    Architecture diagram — horizontal layered boxes.
    Each layer is a colored row with component pills inside.
    """
    headline = data.get("headline", "Reference Architecture")
    _add_header_bar(slide, headline)
    _add_footer(slide, slide_number)

    layers = data.get("layers", [])
    layer_colors = [
        _rgb(0xE3, 0xF2, 0xFD),  # Ingestion — pale blue
        _rgb(0xE8, 0xF5, 0xE9),  # Processing — pale green
        _rgb(0xFF, 0xF9, 0xC4),  # Storage — pale yellow
        _rgb(0xFB, 0xE9, 0xE7),  # AI/ML — pale red
        _rgb(0xF3, 0xE5, 0xF5),  # Visualization — pale purple
    ]
    label_colors = [BLUE, GREEN, _rgb(0xF5, 0x7F, 0x17), RED, PURPLE]

    left = CONTENT_LEFT
    top = CONTENT_TOP + Cm(0.4)
    total_height = CONTENT_HEIGHT - Cm(0.5)
    layer_height = total_height / max(len(layers), 1)

    label_width = Cm(3.2)
    components_left = left + label_width + Cm(0.3)
    components_width = CONTENT_WIDTH - label_width - Cm(0.3)

    for idx, layer in enumerate(layers):
        name = layer.get("name", "")
        components = layer.get("components", [])
        color = layer_colors[idx % len(layer_colors)]
        label_color = label_colors[idx % len(label_colors)]
        layer_top = top + layer_height * idx

        # Layer background
        _add_rect(
            slide,
            left=left,
            top=layer_top,
            width=CONTENT_WIDTH,
            height=layer_height - Cm(0.1),
            fill_color=color,
            line_color=_rgb(0xDD, 0xDD, 0xDD),
            line_width_pt=0.5,
        )

        # Layer label
        _add_textbox(
            slide,
            name,
            left=left + Cm(0.2),
            top=layer_top + (layer_height - Cm(0.8)) / 2,
            width=label_width - Cm(0.3),
            height=Cm(0.8),
            font_size=11,
            bold=True,
            color=label_color,
        )

        # Component pills
        pill_w = (components_width - Cm(0.3) * (len(components) - 1)) / max(len(components), 1)
        pill_h = layer_height - Cm(0.4)

        for comp_idx, comp in enumerate(components):
            pill_left = components_left + (pill_w + Cm(0.3)) * comp_idx
            pill_top = layer_top + Cm(0.2)

            pill = _add_rect(
                slide,
                left=pill_left,
                top=pill_top,
                width=pill_w,
                height=pill_h,
                fill_color=WHITE,
                line_color=label_color,
                line_width_pt=1.0,
            )
            _add_text_to_shape(
                pill,
                comp,
                font_size=11,
                bold=False,
                color=label_color,
                align=PP_ALIGN.CENTER,
            )

    slide.notes_slide.notes_text_frame.text = (
        "Walk the architecture bottom-up or top-down depending on audience. "
        "Technical: go layer by layer — explain each product's role. "
        "Executive: focus on data flow and outcomes, skip product names. "
        "Emphasise managed services = no infrastructure management overhead."
    )


# ---------------------------------------------------------------------------
# Slide type dispatch
# ---------------------------------------------------------------------------

SLIDE_RENDERERS = {
    "title": _render_title,
    "agenda": _render_agenda,
    "pain_points": _render_pain_points,
    "solution_map": _render_solution_map,
    "comparison_table": _render_comparison_table,
    "customer_proof": _render_customer_proof,
    "roi_snapshot": _render_roi_snapshot,
    "timeline": _render_timeline,
    "next_steps": _render_next_steps,
    "text_slide": _render_text_slide,
    "architecture": _render_architecture,
}


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def generate_pitch_deck(config: dict[str, Any]) -> Path:
    """
    Generate a Google Cloud branded pitch deck from a config dictionary.

    Parameters
    ----------
    config:
        Dictionary matching the schema documented in the module docstring.
        Required keys: slides, output_path.
        Optional keys: company_name, industry, presenter, date, template, language.

    Returns
    -------
    Path
        Absolute path to the generated .pptx file.

    Raises
    ------
    ValueError
        If config is missing required keys or contains an unknown slide type.
    FileNotFoundError
        If the output directory does not exist.
    """
    # Validate required fields
    if "slides" not in config:
        raise ValueError("config must contain a 'slides' list")
    if "output_path" not in config:
        raise ValueError("config must contain an 'output_path' string")

    output_path = Path(config["output_path"])
    if not output_path.parent.exists():
        raise FileNotFoundError(f"Output directory does not exist: {output_path.parent}")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout for all slides to retain full control over layout
    blank_layout = prs.slide_layouts[6]  # index 6 is blank in standard template

    slide_number = 0  # title slide gets no number; content slides start at 1

    for slide_data in config["slides"]:
        slide_type = slide_data.get("type", "text_slide")
        renderer = SLIDE_RENDERERS.get(slide_type)

        if renderer is None:
            known = ", ".join(SLIDE_RENDERERS.keys())
            raise ValueError(
                f"Unknown slide type '{slide_type}'. Known types: {known}"
            )

        pptx_slide = prs.slides.add_slide(blank_layout)

        if slide_type == "title":
            renderer(pptx_slide, slide_data, config)
        else:
            slide_number += 1
            renderer(pptx_slide, slide_data, slide_number)

    prs.save(str(output_path))
    return output_path


# ---------------------------------------------------------------------------
# CLI entrypoint
# ---------------------------------------------------------------------------

def _build_sample_config(output_path: str) -> dict[str, Any]:
    """Return the full KGHM sample config for standalone testing."""
    return {
        "company_name": "KGHM Polska Miedz",
        "industry": "Mining & Metals",
        "presenter": "Kacper Ptasinski",
        "date": "2026-03-17",
        "template": "executive-overview",
        "language": "en",
        "output_path": output_path,
        "slides": [
            {
                "type": "title",
                "headline": "KGHM + Google Cloud",
                "subtitle": "Modernizing mining operations with data and AI",
            },
            {
                "type": "agenda",
                "items": [
                    "Your challenges",
                    "Industry context",
                    "Our solution",
                    "Proof points",
                    "Next steps",
                ],
                "timing": ["5 min", "5 min", "15 min", "5 min", "5 min"],
            },
            {
                "type": "pain_points",
                "headline": "Three friction points your teams face daily",
                "points": [
                    {
                        "pain": "Underground safety monitoring is fragmented",
                        "source": "CTO interview, Feb 2026",
                    },
                    {
                        "pain": "Predictive maintenance is reactive, not predictive",
                        "source": "Industry baseline",
                    },
                    {
                        "pain": "NIS2 compliance deadline approaching with no plan",
                        "source": "Regulatory timeline",
                    },
                ],
            },
            {
                "type": "solution_map",
                "headline": "Google Cloud solves all three at once",
                "rows": [
                    {
                        "challenge": "Fragmented safety data",
                        "solution": "BigQuery + Vertex AI",
                        "outcome": "Real-time anomaly detection",
                    },
                    {
                        "challenge": "Reactive maintenance",
                        "solution": "Vertex AI + Pub/Sub",
                        "outcome": "7-day failure prediction",
                    },
                    {
                        "challenge": "NIS2 compliance gap",
                        "solution": "Chronicle SIEM + SCC",
                        "outcome": "24h incident reporting",
                    },
                ],
            },
            {
                "type": "comparison_table",
                "headline": "Google Cloud vs Current State",
                "headers": ["Dimension", "Today", "With Google Cloud"],
                "rows": [
                    ["Data processing", "Batch, T+1", "Real-time streaming"],
                    ["Maintenance", "Calendar-based", "AI-predicted"],
                    ["Security monitoring", "Manual log review", "Automated SIEM/SOAR"],
                ],
            },
            {
                "type": "customer_proof",
                "headline": "Mining companies already on Google Cloud",
                "stories": [
                    {
                        "company": "BHP",
                        "result": "40% reduction in unplanned downtime",
                        "detail": "Vertex AI predictive maintenance",
                    },
                    {
                        "company": "Rio Tinto",
                        "result": "Autonomous haul trucks powered by AI",
                        "detail": "Computer vision + Vertex AI",
                    },
                ],
            },
            {
                "type": "roi_snapshot",
                "headline": "The financial case is straightforward",
                "metrics": [
                    {
                        "label": "Maintenance cost reduction",
                        "value": "30-40%",
                        "detail": "Predictive vs reactive",
                    },
                    {
                        "label": "Energy optimization",
                        "value": "10-15%",
                        "detail": "AI-driven process control",
                    },
                    {
                        "label": "Compliance cost avoidance",
                        "value": "€2-5M/yr",
                        "detail": "NIS2 penalty prevention",
                    },
                ],
            },
            {
                "type": "timeline",
                "headline": "How we get there — 3 phases",
                "phases": [
                    {
                        "name": "Pilot",
                        "duration": "0-6 months",
                        "items": [
                            "1 mine site",
                            "Predictive maintenance POC",
                            "Chronicle deployment",
                        ],
                    },
                    {
                        "name": "Expand",
                        "duration": "6-18 months",
                        "items": [
                            "All Polish mines",
                            "Energy optimization",
                            "Full NIS2 compliance",
                        ],
                    },
                    {
                        "name": "Transform",
                        "duration": "18-36 months",
                        "items": [
                            "Global rollout",
                            "Autonomous operations",
                            "Digital twin",
                        ],
                    },
                ],
            },
            {
                "type": "next_steps",
                "headline": "Concrete next steps",
                "actions": [
                    {
                        "action": "Technical deep-dive with your CTO",
                        "owner": "Google CE",
                        "date": "March 28",
                    },
                    {
                        "action": "Provide query logs for POC sizing",
                        "owner": "Your team",
                        "date": "April 3",
                    },
                    {
                        "action": "POC kickoff",
                        "owner": "Joint",
                        "date": "April 14",
                    },
                ],
            },
            {
                "type": "text_slide",
                "headline": "Why Google Cloud for Mining",
                "body": (
                    "Google Cloud is the only hyperscaler with purpose-built AI infrastructure "
                    "deployed at mining edge sites today. Our Distributed Cloud Edge runs in "
                    "disconnected underground environments — where connectivity is intermittent "
                    "and latency is critical.\n\n"
                    "We bring the same AI platform used in Google Search and DeepMind to "
                    "your ore processing plants, longwall operations, and environmental monitoring."
                ),
            },
            {
                "type": "architecture",
                "headline": "Reference Architecture — Mining Data Platform",
                "layers": [
                    {
                        "name": "Ingestion",
                        "components": ["Pub/Sub", "Datastream", "GDC Edge"],
                    },
                    {
                        "name": "Processing",
                        "components": ["Dataflow", "Dataform"],
                    },
                    {
                        "name": "Storage",
                        "components": ["BigQuery", "Cloud Storage", "Bigtable"],
                    },
                    {
                        "name": "AI/ML",
                        "components": ["Vertex AI", "BigQuery ML"],
                    },
                    {
                        "name": "Visualization",
                        "components": ["Looker", "Looker Studio"],
                    },
                ],
            },
        ],
    }


def main(argv: list[str] | None = None) -> None:
    """CLI entrypoint for pitch_deck.py."""
    parser = argparse.ArgumentParser(
        description="Generate a Google Cloud branded pitch deck (.pptx)",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=(
            "Examples:\n"
            "  python pitch_deck.py                                   # generate sample KGHM deck\n"
            "  python pitch_deck.py --output /tmp/my_deck.pptx        # sample deck to custom path\n"
            "  python pitch_deck.py --config my_config.json           # use custom config\n"
        ),
    )
    parser.add_argument(
        "--config",
        metavar="PATH",
        help="Path to a JSON config file (see module docstring for schema)",
    )
    parser.add_argument(
        "--output",
        metavar="PATH",
        default="/tmp/KGHM_pitch_deck.pptx",
        help="Output .pptx file path (default: /tmp/KGHM_pitch_deck.pptx)",
    )
    args = parser.parse_args(argv)

    if args.config:
        config_path = Path(args.config)
        if not config_path.exists():
            print(f"ERROR: config file not found: {config_path}", file=sys.stderr)
            sys.exit(1)
        with config_path.open() as fh:
            config = json.load(fh)
        # Allow CLI --output to override config output_path
        if args.output != "/tmp/KGHM_pitch_deck.pptx" or "output_path" not in config:
            config["output_path"] = args.output
    else:
        config = _build_sample_config(args.output)

    print(f"Generating pitch deck for: {config.get('company_name', 'Unknown')}")
    print(f"  Slides: {len(config['slides'])}")
    print(f"  Output: {config['output_path']}")

    output = generate_pitch_deck(config)
    print(f"Done. File saved to: {output}")


if __name__ == "__main__":
    main()
